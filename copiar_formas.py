# instalar biblioteca: pip install python-pptx
# se precisar instalar tkinter (reinstalar python): tcl/tk and IDLE
from pptx import Presentation
from copy import deepcopy
from tkinter import Tk
from tkinter.filedialog import askopenfilename
from io import BytesIO

try:
    # Selecionar arquivos
    Tk().withdraw()

    arquivo_origem = askopenfilename(
        title="Selecione a apresentação de origem",
        filetypes=[("PowerPoint", "*.pptx")]
    )

    arquivo_destino = askopenfilename(
        title="Selecione a apresentação de destino",
        filetypes=[("PowerPoint", "*.pptx")]
    )

    # Carrega as apresentações
    apres_origem = Presentation(arquivo_origem)
    apres_destino = Presentation(arquivo_destino)

    # função copiar slides
    def copiar_slide(slide_origem, apresentacao_destino):

       # cria slide em branco
        slide_novo = apresentacao_destino.slides.add_slide(
            apresentacao_destino.slide_layouts[6]
        )

        for shape in slide_origem.shapes:

            # copia placeholders
            if shape.is_placeholder:

                novo_elemento = deepcopy(shape.element)

                slide_novo.shapes._spTree.insert_element_before(
                    novo_elemento,
                    'p:extLst'
                )

                # ajusta posição manualmente
                shape_nova = slide_novo.shapes[-1]

                shape_nova.left = shape.left
                shape_nova.top = shape.top
                shape_nova.width = shape.width
                shape_nova.height = shape.height

            # copia imagens
            elif shape.shape_type == 13:

                imagem = shape.image

                imagem_bytes = BytesIO(imagem.blob)

                slide_novo.shapes.add_picture(
                    imagem_bytes,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )

            # copia outras shapes
            else:

                novo_elemento = deepcopy(shape.element)

                slide_novo.shapes._spTree.insert_element_before(
                    novo_elemento,
                    'p:extLst'
                )

    # função mover slide
    def mover_slide(apresentacao, indice_antigo, indice_novo):

        slides = apresentacao.slides._sldIdLst

        slide = slides[indice_antigo]

        slides.remove(slide)

        slides.insert(indice_novo, slide)

    # copia primeiro slide
    copiar_slide(apres_origem.slides[0], apres_destino)

    # move o slide recém-criado para o início
    mover_slide(
        apres_destino,
        len(apres_destino.slides) - 1,
        0
    )

    # copia último slide
    copiar_slide(apres_origem.slides[-1], apres_destino)

    # Verifica quantidade de slides
    if len(apres_origem.slides) != len(apres_destino.slides):
        raise ValueError(
            "ERRO: As apresentações precisam ter o mesmo número de slides."
        )

    # Percorre os slides (ignora primeiro e ultimo que já foram copiados)
    for i in range(1, len(apres_origem.slides) - 1):
        slide_origem = apres_origem.slides[i]
        slide_destino = apres_destino.slides[i]

        # Copia apenas formas com texto
        for shape in slide_origem.shapes:
            # se a forma tem atributo texto e é diferente de vazio
            if hasattr(shape, "text") and shape.text.strip() != "":
                novo_elemento = deepcopy(shape.element)  # copia o XML da forma

                # _spTree é a estrutura interna do XML que guarda as formas do slide
                slide_destino.shapes._spTree.insert_element_before(
                    novo_elemento, 'p:extLst'
                    # insere novo elemento antes da tag p:extLst se tiver
                    # se não, insere no final da spTree
                )

    # Sobrescreve o destino
    apres_destino.save(arquivo_destino)

    print("Apresentação atualizada com sucesso!")

except Exception as erro:
    print(erro)

input("Pressione ENTER para fechar...")
